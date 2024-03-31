from PptxAPI import  *
import pptx
import os
from pptx.dml.color import RGBColor



def update_text_of_textbox(presentation, slide, text_box_id, new_text):
    slide = presentation.slides[(slide - 1)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Preserve formatting of the first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_color = font.color.rgb
                # Clear existing text and apply new text with preserved formatting
                text_frame.clear()  # Clears all text and formatting
                new_run = text_frame.paragraphs[0].add_run()  # New run in first paragraph
                new_run.text = new_text
                # Reapply formatting
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_color
                return slide
           

def titleMaking(topic,titles,presentaion):
    update_text_of_textbox(presentaion, 1, 1, topic)

    for i in range(len(titles)):               
        update_text_of_textbox(presentaion, 3, i+7, titles[i].split(':')[0].strip())
        # for idx, text in enumerate(list_text_boxes(presentaion,i),1):
        #     print(f"Text box {idx}: {text}")

    for i in range(4,9):               
        update_text_of_textbox(presentaion, i, 1, titles[i-4].split(':')[0].strip())


def flatten_list(nested_list):
    flattened_list = []
    for sublist in nested_list:
        for item in sublist:
            flattened_list.append(item)
    return flattened_list




def ReduceTheLines(content):
    lines =0
    words =len(content)
    currentwords=0
    final=[]
    
    content_text = content.split(":")

    for i in content_text:
        currentwords += len(i)
        lines+=1
        final.append(i)
        if len(final)>=12 or currentwords>=50:
            formatted_content = ''.join(final)
            final_content= formatted_content.strip()
            # print(len(final))
            # print(final_content)
            # print(currentwords)
            return final_content
    return final_content


def ContentMaking(titles,presentaion):
    content = GetListOfContent(titles)
    current_slide_number= 0
    for index in range(len(content)):
        formatted_content = ''.join(content[index])
        if "Content:" in formatted_content:
            # Split the formatted content at "Content:" and take the second part
            content_text = formatted_content.split("Content:")[1]
            final_content= content_text.strip()  # Use strip() to remove leading/trailing whitespace
            formatted_content = ReduceTheLines(final_content)
        else:
            print("No content")
            
        slide =update_text_of_textbox(presentaion, index+4, 2, final_content)
    return presentaion



def mainDo(topic,number_slide):

    presentaion = pptx.Presentation(os.path.join("Templates",f"Template{number_slide}.pptx"))
    title = generate_presentaion_title(topic)
    
    
    topic = makeTitles(title)
    print(topic)
    titles =makeListFormatTitles(generate_slide_title(topic))

    titleMaking(topic,titles,presentaion)
    ContentMaking(titles,presentaion)

    filename = os.path.join("generated_ppt" , f"{topic}_prasentation.pptx")
    presentaion.save(filename)
       

    return get_ppt_download_link(topic),convert_pptx_to_html(filename)


def convert_pptx_to_html(filename):
    # Load the PowerPoint presentation
    presentation = pptx.Presentation(filename)

    # Convert each slide to HTML
    html_content = []
    for slide in presentation.slides:
        html_content.append(slide_to_html(slide))

    # Save the HTML content to a file or return as a string
    # html_file_path = 'presentation.html'
    # with open(html_file_path, 'w') as f:
    #     f.write('\n'.join(html_content))

    # Alternatively, return the HTML content directly
    return '\n'.join(html_content)

    # Return the path to the saved HTML file
    # return FileResponse(html_file_path)

def slide_to_html(slide):
    html_content = []

    # Iterate over shapes in the slide
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            # If the shape contains text, add it to the HTML content
            html_content.append(f"<p>{shape.text}</p>")

    # Join the HTML content into a single string
    slide_html = "".join(html_content)

    # Wrap the slide content in a <div> tag
    return f"<div>{slide_html}</div>"



# ""
# def duplicate_slide(prs, slide):
#     """
#     Duplicate a slide with the same layout and content.
    
#     Parameters:
#         prs (pptx.Presentation): The PowerPoint presentation object.
#         slide (pptx.Slide): The slide to be duplicated.
    
#     Returns:
#         pptx.Slide: The duplicated slide.
#     """
#     # Get the layout of the current slide
#     slide_layout = slide.slide_layout
    
#     # Add a new slide using the same layout
#     duplicated_slide = prs.slides.add_slide(slide_layout)
    
#     # Copy the content of the current slide to the new slide
#     for shape in slide.shapes:
#         new_shape = duplicated_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
#         new_shape.text = shape.text
    
#     return duplicated_slide


# def get_current_slide_number(prs, target_slide):
#     """
#     Get the slide number of a target slide in a PowerPoint presentation.
    
#     Parameters:
#         prs (pptx.Presentation): The PowerPoint presentation object.
#         target_slide (pptx.Slide): The slide whose number you want to find.
    
#     Returns:
#         int: The slide number (1-indexed) if the target slide is found, otherwise returns None.
#     """
#     for idx, slide in enumerate(prs.slides, start=1):
#         if slide == target_slide:
#             return idx
#     return None

# def list_text_boxes(presentaion, slide_num):
#     slide = presentaion.slides[slide_num - 1]
#     text_boxes = []
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text:
#             text_boxes.append(shape.text)
#     return text_boxes

# for idx, text in enumerate(list_text_boxes(presentaion,1),1):
#     print(f"Text box {idx}: {text}")